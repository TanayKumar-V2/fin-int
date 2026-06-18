import logging
import fitz  # PyMuPDF
import docx
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table
from docx.text.paragraph import Paragraph
import pptx
import pandas as pd
from typing import Any

logger = logging.getLogger(__name__)

def parse_pdf(path: str) -> list[dict[str, Any]]:
    """Extract text block by block from a PDF, detecting sections by font size."""
    doc = fitz.open(path)
    blocks = []
    current_section = None
    
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        page_dict = page.get_text("dict")
        
        for block in page_dict.get("blocks", []):
            if block.get("type") == 0:  # text block
                text_content = ""
                max_size = 0
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        text_content += span.get("text", "") + " "
                        span_size = span.get("size", 0)
                        if span_size > max_size:
                            max_size = span_size
                            
                text_content = text_content.strip()
                if not text_content:
                    continue
                    
                # Basic section detection: large text that is relatively short
                if max_size > 14 and len(text_content) < 150:
                    current_section = text_content
                    
                blocks.append({
                    "text": text_content,
                    "page_number": page_num + 1,
                    "section": current_section
                })
                
    return blocks

def parse_docx(path: str) -> list[dict[str, Any]]:
    """Extract paragraphs and tables from DOCX, grouping by headings."""
    doc = docx.Document(path)
    blocks = []
    current_section = None
    
    for child in doc.element.body:
        if isinstance(child, CT_P):
            p = Paragraph(child, doc)
            text = p.text.strip()
            if text:
                if p.style and p.style.name and "Heading" in p.style.name:
                    current_section = text
                blocks.append({
                    "text": text,
                    "page_number": None,
                    "section": current_section
                })
        elif isinstance(child, CT_Tbl):
            t = Table(child, doc)
            table_text = ""
            for row in t.rows:
                row_data = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                table_text += "| " + " | ".join(row_data) + " |\n"
            if table_text.strip():
                blocks.append({
                    "text": table_text.strip(),
                    "page_number": None,
                    "section": current_section
                })
                
    return blocks

def parse_pptx(path: str) -> list[dict[str, Any]]:
    """Extract slide text and notes from PPTX."""
    prs = pptx.Presentation(path)
    blocks = []
    
    for i, slide in enumerate(prs.slides):
        slide_text = ""
        title = None
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text.strip()
            slide_text += f"Title: {title}\n"
            
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                slide_text += shape.text.strip() + "\n"
                
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            
        if notes:
            slide_text += f"\n[Notes]: {notes}"
            
        slide_text = slide_text.strip()
        if slide_text:
            blocks.append({
                "text": slide_text,
                "page_number": i + 1,
                "section": title
            })
            
    return blocks

def parse_xlsx(path: str) -> list[dict[str, Any]]:
    """Extract sheets from XLSX as markdown tables."""
    xl = pd.ExcelFile(path)
    blocks = []
    
    for sheet_name in xl.sheet_names:
        df = xl.parse(sheet_name)
        if df.empty:
            continue
        text = df.to_markdown(index=False)
        blocks.append({
            "text": text,
            "page_number": None,
            "section": sheet_name
        })
        
    return blocks

def parse_document(path: str, ext: str) -> list[dict[str, Any]]:
    """Dispatcher function to parse a document based on its extension."""
    ext = ext.lower()
    try:
        if ext == ".pdf":
            return parse_pdf(path)
        elif ext == ".docx":
            return parse_docx(path)
        elif ext == ".pptx":
            return parse_pptx(path)
        elif ext == ".xlsx":
            return parse_xlsx(path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")
    except Exception as e:
        logger.error(f"Error parsing {path}: {e}")
        raise
